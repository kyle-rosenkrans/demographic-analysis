#!/usr/bin/env python3
"""Extract all text + shape/image info from the InSite pptx."""
from pptx import Presentation
from pptx.util import Emu
import sys, json, os

path = "/Users/kylerosenkrans/Library/CloudStorage/GoogleDrive-krosenkrans@apps.teamschools.org/My Drive/KIPP Miami - Broward Demographics - MEETING EDITS.pptx"
prs = Presentation(path)
out = []
for i, slide in enumerate(prs.slides, 1):
    s = {"slide": i, "texts": [], "tables": [], "images": 0, "shapes": 0}
    for shape in slide.shapes:
        s["shapes"] += 1
        try:
            if shape.has_text_frame:
                txt = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
                if txt:
                    s["texts"].append(txt)
            if shape.has_table:
                tbl = []
                for row in shape.table.rows:
                    tbl.append([c.text.strip() for c in row.cells])
                s["tables"].append(tbl)
            if shape.shape_type == 13:  # picture
                s["images"] += 1
        except Exception as e:
            s["texts"].append(f"[err: {e}]")
    out.append(s)

print(f"TOTAL SLIDES: {len(out)}")
for s in out:
    print("=" * 80)
    print(f"SLIDE {s['slide']} — shapes={s['shapes']} images={s['images']}")
    for t in s["texts"]:
        print(t)
    for tbl in s["tables"]:
        print("--TABLE--")
        for row in tbl:
            print(" | ".join(row))
